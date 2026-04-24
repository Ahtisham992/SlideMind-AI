import io
from typing import List
from models.schemas import SlideContent, ParsedDocument


def parse_pdf(file_bytes: bytes, filename: str) -> ParsedDocument:
    """Parse PDF file and extract text from each page."""
    try:
        import pdfminer.high_level as pdfminer
        from pdfminer.layout import LAParams
        
        slides = []
        raw_text_parts = []
        
        # Extract text page by page
        with io.BytesIO(file_bytes) as f:
            full_text = pdfminer.extract_text(f, laparams=LAParams())
        
        # Split into pages (approximation)
        pages = full_text.split('\x0c')
        pages = [p.strip() for p in pages if p.strip()]
        
        if not pages:
            pages = [full_text]
        
        for i, page_text in enumerate(pages):
            lines = page_text.split('\n')
            title = lines[0].strip() if lines else f"Slide {i+1}"
            content = '\n'.join(lines[1:]).strip() if len(lines) > 1 else page_text
            
            slides.append(SlideContent(
                slide_number=i + 1,
                title=title[:100] if title else f"Slide {i+1}",
                content=content or page_text
            ))
            raw_text_parts.append(page_text)
        
        return ParsedDocument(
            filename=filename,
            total_slides=len(slides),
            slides=slides,
            raw_text='\n\n'.join(raw_text_parts)
        )
    except Exception as e:
        raise ValueError(f"Failed to parse PDF: {str(e)}")


def parse_pptx(file_bytes: bytes, filename: str) -> ParsedDocument:
    """Parse PowerPoint file and extract text from each slide."""
    try:
        from pptx import Presentation
        
        prs = Presentation(io.BytesIO(file_bytes))
        slides = []
        raw_text_parts = []
        
        for i, slide in enumerate(prs.slides):
            title_text = ""
            content_parts = []
            
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                
                text = shape.text_frame.text.strip()
                if not text:
                    continue
                
                # Try to detect title
                if shape.shape_type == 13 or (hasattr(shape, 'placeholder_format') and 
                   shape.placeholder_format is not None and 
                   shape.placeholder_format.idx == 0):
                    title_text = text
                else:
                    content_parts.append(text)
            
            # Fallback: first shape text is title
            if not title_text and content_parts:
                title_text = content_parts.pop(0)
            
            content = '\n'.join(content_parts)
            full_slide_text = f"{title_text}\n{content}".strip()
            
            slides.append(SlideContent(
                slide_number=i + 1,
                title=title_text[:100] if title_text else f"Slide {i+1}",
                content=content
            ))
            raw_text_parts.append(full_slide_text)
        
        return ParsedDocument(
            filename=filename,
            total_slides=len(slides),
            slides=slides,
            raw_text='\n\n'.join(raw_text_parts)
        )
    except Exception as e:
        raise ValueError(f"Failed to parse PPTX: {str(e)}")


def parse_file(file_bytes: bytes, filename: str, content_type: str) -> ParsedDocument:
    """Route file to appropriate parser based on type."""
    filename_lower = filename.lower()
    
    if filename_lower.endswith('.pdf') or 'pdf' in content_type:
        return parse_pdf(file_bytes, filename)
    elif filename_lower.endswith(('.pptx', '.ppt')) or 'presentation' in content_type:
        return parse_pptx(file_bytes, filename)
    else:
        raise ValueError(f"Unsupported file type: {filename}. Please upload PDF or PPTX files.")